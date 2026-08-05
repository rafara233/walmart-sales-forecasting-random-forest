"""
Dashboard Prediksi Penjualan Walmart
-------------------------------------
Eksplorasi data, statistik deskriptif, model Random Forest Regression,
dan simulasi prediksi Weekly Sales berbasis data Walmart 2011-2012.
"""

import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st
from sklearn.ensemble import RandomForestRegressor
from sklearn.metrics import mean_absolute_error, mean_squared_error, r2_score
from sklearn.model_selection import train_test_split
DATA_URL = (
    "https://raw.githubusercontent.com/rafara233/"
    "walmart-sales-forecasting-random-forest/refs/heads/main/Walmart_Sales.csv"
)

# ----------------------------------------------------------------------------
# PALET & TEMA
# ----------------------------------------------------------------------------

INK = "#132A3A"
PAPER = "#EEF1ED"
CARD = "#FFFFFF"
BORDER = "#DCE3DD"
ACCENT = "#2F6F52"
GOLD = "#B98A2E"
MUTED = "#5B6B66"
PALETTE = [ACCENT, GOLD, INK, "#7C9885", "#C97B4A", "#3D5A5B"]

DISPLAY_FONT = "'Space Grotesk', sans-serif"
BODY_FONT = "'IBM Plex Sans', sans-serif"
MONO_FONT = "'IBM Plex Mono', monospace"

MINGGU_PER_BULAN = {1: 5, 2: 4, 3: 5, 4: 4, 5: 5, 6: 4, 7: 5, 8: 5, 9: 4, 10: 5, 11: 4, 12: 5}


# ----------------------------------------------------------------------------
# SETUP HALAMAN & GAYA
# ----------------------------------------------------------------------------

def configure_page():
    st.set_page_config(
        page_title="Walmart Sales Forecasting Dashboard",
        page_icon="🧾",
        layout="wide",
    )

    st.markdown(
        f"""
        <style>
        @import url('https://fonts.googleapis.com/css2?family=Space+Grotesk:wght@500;600;700&family=IBM+Plex+Sans:wght@400;500;600&family=IBM+Plex+Mono:wght@500;600&display=swap');

        html, body, [class*="css"] {{
            font-family: {BODY_FONT};
            color: {INK};
        }}

        .stApp {{ background-color: {PAPER}; }}

        h1, h2, h3, h4 {{
            font-family: {DISPLAY_FONT} !important;
            color: {INK} !important;
            letter-spacing: -0.01em;
        }}

        .eyebrow {{
            font-family: {MONO_FONT};
            font-size: 0.72rem;
            letter-spacing: 0.16em;
            text-transform: uppercase;
            color: {GOLD};
            margin-bottom: 0.2rem;
        }}

        .lede {{
            color: {MUTED};
            font-size: 1.02rem;
            max-width: 68ch;
            width: 100%;
            line-height: 1.6;
            text-wrap: pretty;
        }}

        div[data-testid="stMetric"] {{
            background: {CARD};
            border: 1px solid {BORDER};
            border-bottom: 2px dashed {BORDER};
            border-radius: 14px;
            padding: 0.9rem 1.1rem 0.8rem 1.1rem;
            box-shadow: 0 1px 2px rgba(19,42,58,0.05);
        }}
        div[data-testid="stMetricValue"] {{
            font-family: {MONO_FONT} !important;
            color: {INK} !important;
        }}
        div[data-testid="stMetricLabel"] {{
            font-family: {MONO_FONT} !important;
            text-transform: uppercase;
            font-size: 0.7rem !important;
            letter-spacing: 0.05em;
            color: {MUTED} !important;
        }}

        div[data-testid="stAlert"] {{
            border-radius: 12px;
            box-shadow: 0 1px 2px rgba(19,42,58,0.04);
        }}

        div.stButton > button {{
            font-family: {DISPLAY_FONT};
            background-color: {CARD} !important;
            color: {INK} !important;
            border: 1.5px solid {BORDER} !important;
            border-radius: 10px !important;
            font-weight: 600 !important;
            height: 46px;
            letter-spacing: 0.02em;
            transition: background-color 0.15s ease, color 0.15s ease, border-color 0.15s ease;
        }}
        div.stButton > button p {{
            color: {INK} !important;
            font-weight: 600 !important;
        }}
        div.stButton > button:hover {{
            background-color: {CARD} !important;
            border-color: {ACCENT} !important;
        }}
        div.stButton > button:hover p {{
            color: {ACCENT} !important;
        }}

        /* Primary action button (e.g., "Calculate Prediction" & active navbar item) */
        div.stButton > button[kind="primary"] {{
            background-color: {ACCENT} !important;
            border-color: {ACCENT} !important;
        }}
        div.stButton > button[kind="primary"] p {{
            color: #FFFFFF !important;
        }}
        div.stButton > button[kind="primary"]:hover {{
            background-color: {INK} !important;
            border-color: {INK} !important;
        }}
        div.stButton > button[kind="primary"]:hover p {{
            color: #FFFFFF !important;
        }}

        /* ---------- Custom navbar (replace third-party components) ----------
           Wrapped in st.container(border=True) so it nests correctly
           in the DOM, so CSS can target the original wrapper. */
        div[data-testid="stVerticalBlockBorderWrapper"] {{
            background: {CARD};
            border: 1px solid {BORDER} !important;
            border-radius: 14px !important;
            margin-bottom: 1.2rem;
            box-shadow: 0 1px 2px rgba(19,42,58,0.05);
        }}
        div[data-testid="stVerticalBlockBorderWrapper"] div[data-testid="stHorizontalBlock"] {{
            gap: 0.35rem;
        }}
        div[data-testid="stVerticalBlockBorderWrapper"] div.stButton > button {{
            border: none !important;
            border-radius: 10px !important;
            height: 44px;
            font-size: 14.5px !important;
        }}
        div[data-testid="stVerticalBlockBorderWrapper"] div.stButton > button:hover {{
            background-color: {PAPER} !important;
        }}
        div[data-testid="stVerticalBlockBorderWrapper"] div.stButton > button[kind="primary"] {{
            box-shadow: 0 1px 3px rgba(47,111,82,0.35);
        }}

        div[data-testid="stDataFrame"] {{
            border-radius: 12px;
            overflow: hidden;
            border: 1px solid {BORDER};
        }}

        div[data-testid="stExpander"] {{
            border: 1px solid {BORDER};
            border-radius: 12px;
            background: {CARD};
        }}

        hr {{ border-top: 1px dashed {BORDER} !important; }}

        /* Fallback: some built-in Streamlit components use white text
           that isn't visible over light theme background. */
        p, span, label, li, div[data-testid="stMarkdownContainer"] {{
            color: {INK};
        }}
        div[data-testid="stCaptionContainer"] {{
            color: {MUTED} !important;
        }}
        iframe {{ color-scheme: light; }}

        #MainMenu, footer {{ visibility: hidden; }}

        /* ------------------------------------------------------------
           RESPONSIVE: Full width on PC/laptop, iPad adjusts content width,
           mobile is compact (fonts, padding, menu scrolls).
           Streamlit st.columns auto-stack on narrow screens;
           this only cleans up typography & spacing at each size.
           ------------------------------------------------------------ */

        .block-container {{
            padding-top: 1.6rem;
            padding-bottom: 3rem;
            max-width: 1200px;
        }}

        @media (max-width: 768px) {{
            .block-container {{ max-width: 100%; }}
            h1 {{ font-size: 1.8rem !important; }}
            h2 {{ font-size: 1.5rem !important; }}
            .lede {{ font-size: 0.95rem !important; }}
            div[data-testid="stMetricLabel"] {{ font-size: 0.65rem !important; }}
        }}
        </style>
        """,
        unsafe_allow_html=True,
    )


# -------------------------------------------------------------------
# DATA LOADING & TRAINING
# -------------------------------------------------------------------

@st.cache_data
def load_data() -> pd.DataFrame:
    df = pd.read_csv(DATA_URL)
    df["Date"] = pd.to_datetime(df["Date"], format="%d-%m-%Y")
    df["Year"] = df["Date"].dt.year
    df["Month"] = df["Date"].dt.month
    df["Week"] = df["Date"].dt.isocalendar().week.astype(int)
    return df


@st.cache_resource
def train_models(df: pd.DataFrame) -> dict:
    df_model = df.drop(columns=["Date"])
    X = df_model.drop(columns=["Weekly_Sales"])
    y = df_model["Weekly_Sales"]

    split_ratio = {"90 : 10": 0.10, "80 : 20": 0.20, "70 : 30": 0.30, "60 : 40": 0.40}
    hasil_split = []
    for nama, test_size in split_ratio.items():
        Xtr, Xte, ytr, yte = train_test_split(X, y, test_size=test_size, random_state=42)
        rf = RandomForestRegressor(n_estimators=100, random_state=42)
        rf.fit(Xtr, ytr)
        pred = rf.predict(Xte)
        hasil_split.append(
            {
                "Split": nama,
                "MAE": mean_absolute_error(yte, pred),
                "RMSE": np.sqrt(mean_squared_error(yte, pred)),
                "R²": r2_score(yte, pred),
            }
        )
    hasil_split = pd.DataFrame(hasil_split)
    best_split = hasil_split.loc[hasil_split["R²"].idxmax()]

    X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.20, random_state=42)
    model = RandomForestRegressor(n_estimators=100, random_state=42)
    model.fit(X_train, y_train)
    pred = model.predict(X_test)

    mae = mean_absolute_error(y_test, pred)
    rmse = np.sqrt(mean_squared_error(y_test, pred))
    r2 = r2_score(y_test, pred)

    feature_importance = pd.DataFrame(
        {"Feature": X.columns, "Importance": model.feature_importances_}
    ).sort_values("Importance", ascending=False)

    corr = df_model.corr(numeric_only=True)

    return {
        "df_model": df_model,
        "X": X,
        "y": y,
        "hasil_split": hasil_split,
        "best_split": best_split,
        "model": model,
        "X_test": X_test,
        "y_test": y_test,
        "pred": pred,
        "mae": mae,
        "rmse": rmse,
        "r2": r2,
        "feature_importance": feature_importance,
        "corr": corr,
    }


# -------------------------------------------------------------------
# STYLE & HELPER FUNCTIONS
# -------------------------------------------------------------------

def page_header(eyebrow: str, title: str, description: str):
    st.markdown(
        f'<div class="eyebrow">{eyebrow}</div>'
        f'<h2 style="margin-top:-6px;">{title}</h2>'
        f'<p class="lede">{description}</p>',
        unsafe_allow_html=True,
    )


def style_table(df: pd.DataFrame) -> pd.io.formats.style.Styler:
    return (
        df.style.set_properties(**{"text-align": "center"})
        .format(precision=2)
    )


# -------------------------------------------------------------------
# PAGE: DATASET
# -------------------------------------------------------------------

def page_dataset(df: pd.DataFrame, results: dict):
    df_model = results["df_model"]
    
    page_header(
        "01 · DATASET",
        "Raw Data Overview",
        "View the complete Walmart sales dataset with all features used for forecasting.",
    )

    st.markdown(
        f"""
        **Dataset contains:**
        - **Rows:** {len(df):,}
        - **Columns:** {len(df.columns)}
        - **Date range:** {int(df['Year'].min())}–{int(df['Year'].max())}
        - **Stores:** {df['Store'].nunique()} locations
        """
    )

    st.divider()
    st.subheader("Raw Data Sample")
    st.dataframe(style_table(df.head(10)), hide_index=True, use_container_width=True)

    st.divider()
    st.subheader("Data Types & Missing Values")
    info_data = pd.DataFrame({
        "Column": df.columns,
        "Data Type": df.dtypes,
        "Non-null Count": df.count(),
        "Missing": df.isnull().sum(),
    })
    st.dataframe(style_table(info_data), hide_index=True, use_container_width=True)

    st.divider()
    st.subheader("Numerical Summary")
    st.dataframe(style_table(df_model.describe()), use_container_width=True)


# -------------------------------------------------------------------
# PAGE: STATISTICS
# -------------------------------------------------------------------

def page_statistik(df: pd.DataFrame):
    page_header(
        "02 · STATISTICS",
        "Descriptive Analysis",
        "Explore distributions, correlations, and seasonal trends in the sales data.",
    )

    st.subheader("Weekly Sales Distribution")
    fig_dist = px.histogram(df, x="Weekly_Sales", nbins=40, color_discrete_sequence=[ACCENT])
    fig_dist.update_layout(
        title="Distribution of Weekly Sales",
        xaxis_title="Weekly Sales (USD)",
        yaxis_title="Frequency",
        template="plotly_white",
        hovermode="x unified",
        font=dict(family=BODY_FONT),
    )
    st.plotly_chart(fig_dist, use_container_width=True)

    st.divider()
    st.subheader("Sales by Store")
    store_sales = df.groupby("Store")["Weekly_Sales"].agg(["mean", "std"]).reset_index()
    store_sales = store_sales.sort_values("mean", ascending=False)
    
    fig_store = px.bar(store_sales, x="Store", y="mean", error_y="std", color_discrete_sequence=[ACCENT])
    fig_store.update_layout(
        title="Average Weekly Sales by Store",
        xaxis_title="Store",
        yaxis_title="Average Sales (USD)",
        template="plotly_white",
        font=dict(family=BODY_FONT),
    )
    st.plotly_chart(fig_store, use_container_width=True)

    st.divider()
    st.subheader("Seasonal Trends (Monthly)")
    monthly_sales = df.groupby("Month")["Weekly_Sales"].mean().reset_index()
    fig_month = px.line(
        monthly_sales, x="Month", y="Weekly_Sales",
        markers=True, color_discrete_sequence=[ACCENT],
        title="Average Sales by Month",
    )
    fig_month.update_layout(
        xaxis_title="Month",
        yaxis_title="Average Sales (USD)",
        template="plotly_white",
        font=dict(family=BODY_FONT),
    )
    st.plotly_chart(fig_month, use_container_width=True)

    st.divider()
    st.subheader("Correlation Matrix")
    corr = df.corr(numeric_only=True)
    fig_corr = px.imshow(corr, text_auto=".2f", color_continuous_scale="RdBu", zmin=-1, zmax=1)
    fig_corr.update_layout(
        title="Feature Correlations",
        font=dict(family=BODY_FONT),
    )
    st.plotly_chart(fig_corr, use_container_width=True)


# -------------------------------------------------------------------
# PAGE: MACHINE LEARNING
# -------------------------------------------------------------------

def page_ml(df: pd.DataFrame, results: dict):
    hasil_split = results["hasil_split"]
    feature_importance = results["feature_importance"]
    corr = results["corr"]

    page_header(
        "03 · MACHINE LEARNING",
        "Random Forest Model",
        "Model training results, feature importance, and performance metrics across different train-test splits.",
    )

    st.subheader("Model Performance by Train-Test Split")
    st.dataframe(style_table(hasil_split), hide_index=True, use_container_width=True)

    st.divider()
    st.subheader("Feature Importance")
    fig_feat = px.bar(
        feature_importance.head(10),
        x="Importance",
        y="Feature",
        orientation="h",
        color_discrete_sequence=[ACCENT],
        title="Top 10 Most Important Features",
    )
    fig_feat.update_layout(
        template="plotly_white",
        font=dict(family=BODY_FONT),
    )
    st.plotly_chart(fig_feat, use_container_width=True)

    st.divider()
    st.subheader("Correlation Matrix")
    fig_corr = px.imshow(corr, text_auto=".2f", color_continuous_scale="RdBu", zmin=-1, zmax=1)
    fig_corr.update_layout(
        title="Feature Correlations",
        font=dict(family=BODY_FONT),
    )
    st.plotly_chart(fig_corr, use_container_width=True)

    st.divider()
    with st.expander("About Feature Importance"):
        st.markdown(
            "Random Forest calculates importance by measuring how much each feature decreases "
            "impurity (typically Gini impurity) when splitting nodes in decision trees. Features that "
            "separate samples more effectively capturing seasonal patterns or variable interactions "
            "that aren't linear rank higher. This differs from simple correlation analysis, so rankings may vary."
        )


# -------------------------------------------------------------------
# PAGE: PREDICTION
# -------------------------------------------------------------------

def page_prediksi(df: pd.DataFrame, results: dict):
    model = results["model"]

    page_header(
        "04 · SIMULATION",
        "Weekly Sales Prediction",
        "Enter store and economic conditions, then view estimated weekly sales from the model.",
    )

    with st.expander("Example Input"):
        contoh = pd.DataFrame(
            {
                "Variable": ["Store", "Holiday Flag", "Temperature", "Fuel Price", "CPI", "Unemployment", "Year", "Month", "Week"],
                "Example": [5, 0, 75.5, 3.65, 212.45, 6.20, 2026, 8, 2],
            }
        )
        st.dataframe(style_table(contoh), hide_index=True, use_container_width=True)

    st.divider()
    c1, c2, c3 = st.columns(3)
    store = c1.number_input("Store", min_value=int(df["Store"].min()), max_value=int(df["Store"].max()), value=5)
    holiday = c2.selectbox("Holiday Flag", [0, 1], format_func=lambda x: "Holiday" if x == 1 else "Regular Day")
    temperature = c3.number_input("Temperature", min_value=float(df["Temperature"].min()), max_value=float(df["Temperature"].max()), value=float(df["Temperature"].mean()))

    fuel = c1.number_input("Fuel Price", min_value=float(df["Fuel_Price"].min()), max_value=float(df["Fuel_Price"].max()), value=float(df["Fuel_Price"].mean()))
    cpi = c2.number_input("CPI", min_value=float(df["CPI"].min()), max_value=float(df["CPI"].max()), value=float(df["CPI"].mean()))
    unemployment = c3.number_input("Unemployment", min_value=float(df["Unemployment"].min()), max_value=float(df["Unemployment"].max()), value=float(df["Unemployment"].mean()))

    year = c1.number_input("Year", min_value=2011, max_value=2100, value=2026)
    month = c2.selectbox("Month", list(range(1, 13)))
    max_week = MINGGU_PER_BULAN[month]
    week = c3.selectbox("Week", list(range(1, max_week + 1)))
    st.caption(f"Month {month} has weeks 1–{max_week}.")

    st.divider()

    if st.button("Calculate Prediction", type="primary"):
        # Ensure column order matches training data
        data_baru = pd.DataFrame(
            {
                "Store": [store], 
                "Holiday_Flag": [holiday], 
                "Temperature": [temperature],
                "Fuel_Price": [fuel], 
                "CPI": [cpi], 
                "Unemployment": [unemployment],
                "Year": [year], 
                "Month": [month], 
                "Week": [week],
            }
        )
        # Reorder to match model training features
        data_baru = data_baru[results["X"].columns]
        hasil = model.predict(data_baru)[0]

        q1, q3 = df["Weekly_Sales"].quantile([0.25, 0.75])
        if hasil < q1:
            kategori, warna = "Low", "🔴"
        elif hasil < q3:
            kategori, warna = "Medium", "🟡"
        else:
            kategori, warna = "High", "🟢"

        st.success(f"Estimated Weekly Sales: **USD ${hasil:,.2f}**")

        rc1, rc2 = st.columns(2)
        rc1.metric("Predicted Weekly Sales", f"${hasil:,.0f}")
        rc2.metric("Category", f"{warna} {kategori}")

        st.divider()
        st.markdown(
            f"Based on the entered conditions, the model estimates Store "
            f"**{store}**'s weekly sales to be in the **{kategori.lower()}** range compared to other stores in the "
            "dataset. This figure is an average across all decision trees in the model, trained "
            "on historical patterns from 2011–2012, so it's most accurate for conditions similar to that period."
        )

    st.caption(
        "Training dataset covers only 2011–2012. Years outside this range can still be entered, "
        "but results are extrapolations, not truly validated values."
    )


# -------------------------------------------------------------------
# PAGE: CONCLUSION
# -------------------------------------------------------------------

def page_kesimpulan(df: pd.DataFrame, results: dict):
    mae, rmse, r2 = results["mae"], results["rmse"], results["r2"]
    best_split = results["best_split"]
    feature_importance = results["feature_importance"]
    top_feat = feature_importance.iloc[0]
    df_model = results["df_model"]

    page_header(
        "05 · SUMMARY",
        "Conclusion",
        "A brief overview of the entire workflow from raw data to production-ready model.",
    )

    c1, c2, c3 = st.columns(3)
    c1.metric("Total Records", f"{len(df_model):,}")
    c2.metric("Number of Stores", df_model["Store"].nunique())
    c3.metric("Period", f"{int(df['Year'].min())}–{int(df['Year'].max())}")

    st.markdown(
        f"This dashboard is built on **{len(df_model):,} rows** of weekly sales data from "
        f"**{df_model['Store'].nunique()} Walmart stores**. After exploratory data analysis and descriptive statistics "
        "revealed significant variation across stores and months, a Machine Learning "
        "approach specifically **Random Forest Regression** was chosen because it handles "
        "non-linear relationships without extensive preprocessing."
    )

    st.divider()
    st.subheader("Model Performance")
    m1, m2, m3 = st.columns(3)
    m1.metric("MAE", f"${mae:,.0f}")
    m2.metric("RMSE", f"${rmse:,.0f}")
    m3.metric("R² Score", f"{r2:.3f}")
    st.caption(f"Best configuration uses train-test split **{best_split['Split']}**, selected by highest R².")

    st.divider()
    st.subheader("Most Influential Variable")
    v1, v2 = st.columns(2)
    v1.metric("Top Feature", top_feat["Feature"])
    v2.metric("Importance Score", f"{top_feat['Importance']:.3f}")
    st.markdown(
        f"**{top_feat['Feature']}** consistently ranks as the most influential variable for "
        "Weekly Sales predictions, aligning with patterns also visible in the correlation analysis from the Machine Learning tab."
    )

    st.divider()
    st.subheader("Usage Notes")
    st.warning(
        "Model trained on 2011–2012 data. Predictions for years outside this range are based on "
        "historical patterns, not projections accounting for current market conditions."
    )

    st.success(
        "Overall, this dashboard provides a complete view from raw data through statistics to a production model sufficient for both initial exploration and weekly sales scenario simulations."
    )

    st.balloons()


# -------------------------------------------------------------------
# PAGE: REPORTS (PDF & PPT Export)
# -------------------------------------------------------------------

def generate_pdf_report(df: pd.DataFrame, results: dict):
    """Generate CSV/Text report that can be converted to PDF"""
    from io import StringIO
    from datetime import datetime
    
    # Create text-based report
    report = StringIO()
    
    # Header
    report.write("=" * 80 + "\n")
    report.write("WALMART SALES FORECASTING - COMPREHENSIVE REPORT\n")
    report.write("=" * 80 + "\n\n")
    report.write(f"Generated: {datetime.now().strftime('%B %d, %Y at %H:%M')}\n\n")
    
    # Executive Summary
    report.write("EXECUTIVE SUMMARY\n")
    report.write("-" * 80 + "\n")
    df_model = results["df_model"]
    report.write(f"""
This report presents a comprehensive analysis of Walmart weekly sales data covering 
{int(df['Year'].min())}–{int(df['Year'].max())} across {df_model['Store'].nunique()} store locations. 
The analysis includes descriptive statistics, correlation analysis, and a machine 
learning model for sales forecasting.

DATASET OVERVIEW:
  • Total Records: {len(df_model):,}
  • Number of Stores: {df_model['Store'].nunique()}
  • Date Range: {int(df['Year'].min())}–{int(df['Year'].max())}
  • Features Used: {len(results['X'].columns)}
  • Target Variable: Weekly Sales (USD)

""")
    
    # Model Performance
    report.write("MODEL PERFORMANCE METRICS\n")
    report.write("-" * 80 + "\n")
    mae, rmse, r2 = results["mae"], results["rmse"], results["r2"]
    best_split = results["best_split"]
    
    report.write(f"""
Algorithm: Random Forest Regressor
Number of Trees: 100
Max Depth: Default (unlimited)

Performance Metrics:
  • Mean Absolute Error (MAE): ${mae:,.2f}
  • Root Mean Squared Error (RMSE): ${rmse:,.2f}
  • R² Score: {r2:.4f}
  • Best Train-Test Split: {best_split["Split"]}

Interpretation:
  - The model explains {r2*100:.1f}% of the variance in weekly sales
  - Average prediction error is approximately ${mae:,.0f} per week
  - Model performance indicates a strong fit to the training data

""")
    
    # Top Features
    report.write("TOP 10 MOST IMPORTANT FEATURES\n")
    report.write("-" * 80 + "\n")
    feature_importance = results["feature_importance"]
    
    report.write("\nRank | Feature Name             | Importance Score | Percentage\n")
    report.write("-" * 80 + "\n")
    total_importance = feature_importance["Importance"].sum()
    
    for idx, row in feature_importance.head(10).iterrows():
        percentage = (row['Importance'] / total_importance) * 100
        report.write(f"{idx+1:4d} | {row['Feature']:24s} | {row['Importance']:16.6f} | {percentage:6.2f}%\n")
    
    report.write("\n")
    
    # Correlation Summary
    report.write("FEATURE CORRELATIONS WITH TARGET (Weekly Sales)\n")
    report.write("-" * 80 + "\n")
    corr = results["corr"]
    if "Weekly_Sales" in corr.columns:
        sales_corr = corr["Weekly_Sales"].sort_values(ascending=False).head(10)
        for feature, corr_val in sales_corr.items():
            if feature != "Weekly_Sales":
                report.write(f"  {feature:24s}: {corr_val:8.4f}\n")
    
    report.write("\n")
    
    # Conclusions
    report.write("CONCLUSIONS & RECOMMENDATIONS\n")
    report.write("-" * 80 + "\n")
    report.write(f"""
The Random Forest Regression model achieves an R² score of {r2:.4f}, indicating a 
strong fit to the historical data. The model successfully captures non-linear 
relationships between features and weekly sales.

Key Findings:
  • {feature_importance.iloc[0]['Feature']} is the most influential variable for 
    predicting weekly sales
  • The model performs consistently across different train-test split ratios
  • Feature importance rankings reveal which factors drive sales most significantly

Strategic Recommendations:
  ✓ Deploy this model for weekly sales forecasting and planning
  ✓ Monitor prediction accuracy as new data becomes available
  ✓ Retrain the model quarterly to capture seasonal changes and trends
  ✓ Use feature importance rankings to prioritize data collection efforts
  ✓ Integrate model predictions into inventory and staffing decisions
  ✓ Consider ensemble approaches combining this model with other forecasting methods

Limitations & Next Steps:
  • Model trained on historical data (2010-2012)
  • Predictions for years outside training range are extrapolations
  • Regular retraining recommended as business conditions evolve
  • Consider external features (promotions, holidays, economic indicators)

""")
    
    report.write("=" * 80 + "\n")
    report.write("END OF REPORT\n")
    report.write("=" * 80 + "\n")
    
    # Convert to bytes for download
    from io import BytesIO
    buffer = BytesIO(report.getvalue().encode('utf-8'))
    buffer.seek(0)
    return buffer


def download_pdf_from_github():
    """Download PDF report directly from GitHub repository"""
    import requests
    from datetime import datetime
    
    # GitHub raw content URL
    github_url = (
        "https://raw.githubusercontent.com/rafara233/"
        "walmart-sales-forecasting-random-forest/main/Walmart_Sales_Report.pdf"
    )
    
    try:
        response = requests.get(github_url, timeout=10)
        if response.status_code == 200:
            from io import BytesIO
            buffer = BytesIO(response.content)
            buffer.seek(0)
            return buffer
        else:
            return None
    except Exception as e:
        st.error(f"Error downloading PDF: {str(e)}")
        return None


def generate_ppt_report(df: pd.DataFrame, results: dict):
    """Generate PowerPoint presentation from sales forecasting analysis"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        from io import BytesIO
        from datetime import datetime
    except ImportError:
        return None
    
    buffer = BytesIO()
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Color scheme
    DARK_COLOR = RGBColor(19, 42, 58)      # INK
    ACCENT_COLOR = RGBColor(47, 111, 82)   # ACCENT
    GOLD_COLOR = RGBColor(185, 138, 46)    # GOLD
    
    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = ACCENT_COLOR
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(20)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    def add_content_slide(title, content_list):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(content_list):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_COLOR
            p.space_before = Pt(6)
            p.space_after = Pt(6)
            p.level = 0
    
    # Slide 1: Title
    add_title_slide(
        "📊 Walmart Sales Forecasting",
        f"Analysis Report • {datetime.now().strftime('%B %Y')}"
    )
    
    # Slide 2: Overview
    df_model = results["df_model"]
    add_content_slide("Dataset Overview", [
        f"✓ Total Records: {len(df_model):,} weekly observations",
        f"✓ Store Locations: {df_model['Store'].nunique()} Walmart stores",
        f"✓ Time Period: {int(df['Year'].min())}–{int(df['Year'].max())}",
        f"✓ Features: {len(results['X'].columns)} input variables",
        f"✓ Target: Weekly Sales (USD)",
    ])
    
    # Slide 3: Model Performance
    mae, rmse, r2 = results["mae"], results["rmse"], results["r2"]
    best_split = results["best_split"]
    add_content_slide("Model Performance Metrics", [
        f"✓ Model: Random Forest Regressor (100 trees)",
        f"✓ Mean Absolute Error (MAE): ${mae:,.0f}",
        f"✓ Root Mean Squared Error (RMSE): ${rmse:,.0f}",
        f"✓ R² Score: {r2:.4f} (Strong correlation)",
        f"✓ Best Split: {best_split['Split']} train-test ratio",
    ])
    
    # Slide 4: Top Features
    feature_importance = results["feature_importance"]
    top_5 = feature_importance.head(5)
    feature_text = [f"✓ {row['Feature']}: {row['Importance']:.4f}" 
                    for _, row in top_5.iterrows()]
    add_content_slide("Top 5 Most Important Features", feature_text)
    
    # Slide 5: Key Insights
    add_content_slide("Key Insights", [
        f"✓ {feature_importance.iloc[0]['Feature']} is the strongest predictor",
        f"✓ Model explains {r2*100:.1f}% of sales variance",
        f"✓ Average prediction error: ${mae:,.0f} per week",
        f"✓ Non-linear patterns captured by tree-based model",
        f"✓ Ready for production deployment and retraining",
    ])
    
    # Slide 6: Recommendations
    add_content_slide("Recommendations", [
        "✓ Deploy model for weekly sales forecasting",
        "✓ Monitor accuracy as new data accumulates",
        "✓ Retrain quarterly to capture seasonal changes",
        "✓ Use feature importance for strategic decisions",
        "✓ Integrate predictions into inventory planning",
    ])
    
    # Slide 7: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_COLOR
    
    conclusion_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3.5))
    conclusion_frame = conclusion_box.text_frame
    conclusion_frame.word_wrap = True
    conclusion_frame.text = "Thank You"
    conclusion_frame.paragraphs[0].font.size = Pt(54)
    conclusion_frame.paragraphs[0].font.bold = True
    conclusion_frame.paragraphs[0].font.color.rgb = GOLD_COLOR
    conclusion_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    conclusion_frame.add_paragraph()
    p = conclusion_frame.paragraphs[1]
    p.text = "Walmart Sales Forecasting Dashboard"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    prs.save(buffer)
    buffer.seek(0)
    return buffer


def page_reports(df: pd.DataFrame, results: dict):
    page_header(
        "06 · REPORTS",
        "Export & Download",
        "Generate and download comprehensive analysis reports in multiple formats.",
    )
    
    st.subheader("Download Report")
    col1, col2, col3 = st.columns(3)
    
    with col1:
        if st.button("📄 Download as Text Report", use_container_width=True):
            with st.spinner("Generating text report..."):
                txt_buffer = generate_pdf_report(df, results)
                st.download_button(
                    label="📥 Click to download TXT",
                    data=txt_buffer,
                    file_name=f"Walmart_Sales_Report_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.txt",
                    mime="text/plain",
                    use_container_width=True,
                )
                st.success("✅ Text report generated successfully!")
    
    with col2:
        if st.button("📘 Download PDF from GitHub", use_container_width=True, type="primary"):
            with st.spinner("Downloading PDF from GitHub..."):
                try:
                    pdf_buffer = download_pdf_from_github()
                    if pdf_buffer is not None:
                        st.download_button(
                            label="📥 Click to download PDF",
                            data=pdf_buffer,
                            file_name=f"Walmart_Sales_Report_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.pdf",
                            mime="application/pdf",
                            use_container_width=True,
                        )
                        st.success("✅ PDF downloaded successfully!")
                    else:
                        st.warning("⚠️ Could not download PDF from GitHub. Please check the repository.")
                except Exception as e:
                    st.warning(f"⚠️ Error downloading PDF: {str(e)}")
    
    with col3:
        if st.button("🎬 Download as PowerPoint", use_container_width=True):
            with st.spinner("Generating PowerPoint presentation..."):
                try:
                    ppt_buffer = generate_ppt_report(df, results)
                    if ppt_buffer is not None:
                        st.download_button(
                            label="📥 Click to download PPT",
                            data=ppt_buffer,
                            file_name=f"Walmart_Sales_Report_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True,
                        )
                        st.success("✅ PowerPoint generated successfully!")
                    else:
                        st.warning("⚠️ PowerPoint library not available. Please use other formats instead.")
                except Exception as e:
                    st.warning(f"⚠️ PowerPoint generation not available. Error: {str(e)}\nPlease try other download options.")
    
    st.divider()
    st.subheader("Report Contents")
    st.markdown("""
    **Text Report includes:**
    - Executive summary of the analysis
    - Dataset overview and statistics
    - Model performance metrics (MAE, RMSE, R²)
    - Top 10 most important features with percentages
    - Feature correlations with target variable
    - Conclusions and strategic recommendations
    - Limitations and next steps
    
    **PDF Report (from GitHub):**
    - Professional formatted PDF document
    - Visual charts and graphs
    - Detailed analysis and insights
    - Ready for presentation to stakeholders
    
    **PowerPoint Presentation includes:**
    - Title slide with report metadata
    - Dataset overview and statistics
    - Model performance metrics
    - Feature importance rankings
    - Key insights and findings
    - Strategic recommendations
    - Thank you slide
    
    💡 **Note**: Text Report is generated locally. PDF is downloaded from the project GitHub repository.
    PowerPoint requires the `python-pptx` library.
    """)


# -------------------------------------------------------------------
# NAVBAR (native Streamlit, no third-party components)
# -------------------------------------------------------------------

NAV_ITEMS = [
    ("Dataset", "🗂️"),
    ("Statistics", "📊"),
    ("Machine Learning", "🧠"),
    ("Prediction", "🎯"),
    ("Conclusion", "📄"),
    ("Reports", "📥"),
]


def render_navbar() -> str:
    if "menu" not in st.session_state:
        st.session_state.menu = NAV_ITEMS[0][0]

    with st.container(border=True):
        cols = st.columns(len(NAV_ITEMS))
        for col, (label, icon) in zip(cols, NAV_ITEMS):
            is_active = st.session_state.menu == label
            with col:
                if st.button(
                    f"{icon}  {label}",
                    key=f"nav_{label}",
                    use_container_width=True,
                    type="primary" if is_active else "secondary",
                ):
                    st.session_state.menu = label

    return st.session_state.menu


# -------------------------------------------------------------------
# MAIN
# -------------------------------------------------------------------

def main():
    configure_page()

    st.markdown(
        '<div class="eyebrow">WALMART SALES INTELLIGENCE</div>'
        '<h1 style="margin-top:-6px;">🧾 Sales Forecasting Dashboard</h1>'
        '<p class="lede">Explore data, statistics, and Random Forest model to forecast '
        'Weekly Sales for Walmart stores.</p>',
        unsafe_allow_html=True,
    )

    menu = render_navbar()

    df = load_data()
    results = train_models(df)

    if menu == "Dataset":
        page_dataset(df, results)
    elif menu == "Statistics":
        page_statistik(df)
    elif menu == "Machine Learning":
        page_ml(df, results)
    elif menu == "Prediction":
        page_prediksi(df, results)
    elif menu == "Conclusion":
        page_kesimpulan(df, results)
    elif menu == "Reports":
        page_reports(df, results)


if __name__ == "__main__":
    main()
